from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import os
from tkinter import Tk
from tkinter.filedialog import askdirectory, asksaveasfilename

# Hide Tkinter root window
Tk().withdraw()

# 1️⃣ Select folder containing images
image_folder = askdirectory(title="Select Folder Containing Images")
if not image_folder:
    print("No folder selected. Exiting...")
    exit()

# 2️⃣ Select where to save PPTX
output_ppt = asksaveasfilename(
    title="Save PowerPoint As",
    defaultextension=".pptx",
    filetypes=[("PowerPoint Files", "*.pptx")]
)
if not output_ppt:
    print("No PPTX file selected. Exiting...")
    exit()

# 3️⃣ Optional: Select PDF save location
export_pdf = asksaveasfilename(
    title="Save PDF As (Optional)",
    defaultextension=".pdf",
    filetypes=[("PDF Files", "*.pdf")]
)
if export_pdf == "":
    export_pdf = None

# 4️⃣ Optional: Select folder for JPG slides
output_jpg_folder = askdirectory(title="Select Folder to Save JPG Slides (Optional)")
if output_jpg_folder:
    os.makedirs(output_jpg_folder, exist_ok=True)
else:
    output_jpg_folder = None

# 5️⃣ Create Presentation
prs = Presentation()
prs.slide_width = Inches(8.5)
prs.slide_height = Inches(11)

# 6️⃣ Get all images
image_files = sorted([
    f for f in os.listdir(image_folder)
    if f.lower().endswith((".png", ".jpg", ".jpeg"))
])

if not image_files:
    print("No images found in the folder. Exiting...")
    exit()

# 7️⃣ Create slides
for idx, filename in enumerate(image_files):
    img_path = os.path.join(image_folder, filename)

    # Add blank slide
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Add border rectangle (7.5 x 10 inches centered)
    left = Inches((8.5 - 7.5)/2)
    top = Inches((11 - 10)/2)
    width = Inches(7.5)
    height = Inches(10)
    shape = slide.shapes.add_shape(1, left, top, width, height)  # Rectangle
    shape.line.color.rgb = RGBColor(0,0,0)
    shape.line.width = Pt(3)
    shape.fill.background()  # Transparent fill

    # Insert image inside border
    img_left = left + Inches(0.25)
    img_top = top + Inches(0.25)
    img_width = width - Inches(0.5)
    img_height = height - Inches(0.5)
    slide.shapes.add_picture(img_path, img_left, img_top, img_width, img_height)

    print(f"Added slide {idx+1}: {filename}")

# 8️⃣ Save PPTX
prs.save(output_ppt)
print(f"\n✅ PowerPoint saved: {output_ppt}")

# 9️⃣ PDF and JPG export (Windows + PowerPoint)
try:
    import win32com.client

    ppt_app = win32com.client.Dispatch("PowerPoint.Application")
    presentation = ppt_app.Presentations.Open(os.path.abspath(output_ppt))

    # PDF export
    if export_pdf:
        presentation.SaveAs(os.path.abspath(export_pdf), 32)  # 32 = PDF
        print(f"✅ PDF exported: {export_pdf}")

    # JPG export
    if output_jpg_folder:
        for i, slide in enumerate(presentation.Slides, start=1):
            jpg_path = os.path.join(output_jpg_folder, f"Slide_{i}.jpg")
            jpg_path = os.path.abspath(jpg_path)
            slide.Export(jpg_path, "JPG")
        print(f"✅ Slides exported as JPG: {output_jpg_folder}")

    presentation.Close()
    ppt_app.Quit()

except Exception as e:
    print("⚠️ PDF/JPG export failed:", e)
